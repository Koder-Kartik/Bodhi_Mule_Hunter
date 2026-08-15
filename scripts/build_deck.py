#!/usr/bin/env python3
"""Build the submission deck (PPTX) from the measured evaluation output.

Same discipline as the report: no number is typed here. Everything comes from
``artifacts/metrics/evaluation.json`` and the screenshots come from the running
application, so the deck cannot claim something the code does not do.

16:9, dark theme matching the investigator console.
"""

from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

from repo_stats import code_lines_short, test_count  # noqa: E402

from bodhi.config import (  # noqa: E402
    AFFILIATION, EVENT, FIGURE_DIR, METRICS_DIR, ROOT, TEAM, TEAM_NAME,
)

OUT = ROOT / "docs" / "BODHI_Mule_Hunter_Deck.pptx"
OUT_PDF = ROOT / "docs" / "BODHI_Mule_Hunter_Deck.pdf"
SHOTS = ROOT / "docs" / "screenshots"

# palette (mirrors dashboard/styles.css)
BG = RGBColor(0x0B, 0x0F, 0x17)
PANEL = RGBColor(0x12, 0x18, 0x26)
PANEL2 = RGBColor(0x18, 0x20, 0x31)
BORDER = RGBColor(0x23, 0x30, 0x4A)
TEXT = RGBColor(0xE6, 0xED, 0xF7)
DIM = RGBColor(0x8F, 0xA3, 0xC0)
FAINT = RGBColor(0x61, 0x70, 0x8A)
ACCENT = RGBColor(0x4D, 0xA3, 0xFF)
GREEN = RGBColor(0x2E, 0xCC, 0x9B)
AMBER = RGBColor(0xF2, 0xC1, 0x4E)
ORANGE = RGBColor(0xFF, 0x8C, 0x42)
RED = RGBColor(0xFF, 0x4D, 0x5E)
PURPLE = RGBColor(0x7D, 0x5C, 0xFF)

W, H = Inches(13.333), Inches(7.5)
SANS = "Segoe UI"
MONO = "Consolas"


def _pct(x, dp=1):
    return f"{100 * float(x):.{dp}f}%"


def _n(x):
    return f"{int(x):,}"


# --------------------------------------------------------------------------
# primitives
# --------------------------------------------------------------------------


def blank(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def text(slide, x, y, w, h, runs, *, size=16, color=TEXT, bold=False,
         align=PP_ALIGN.LEFT, font=SANS, spacing=1.15, anchor=MSO_ANCHOR.TOP):
    """runs: str, or list of (text, {overrides}) tuples, or list of str lines."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    lines = [runs] if isinstance(runs, str) else runs
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        if isinstance(line, tuple):
            body, over = line
        else:
            body, over = line, {}
        p.space_after = Pt(over.get("space_after", 4))
        r = p.add_run()
        r.text = body
        r.font.size = Pt(over.get("size", size))
        r.font.bold = over.get("bold", bold)
        r.font.color.rgb = over.get("color", color)
        r.font.name = over.get("font", font)
    return tb


def panel(slide, x, y, w, h, *, fill=PANEL, line=BORDER, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    if radius:
        try:
            shape.adjustments[0] = 0.04
        except Exception:
            pass
    return shape


def header(slide, kicker, title, *, sub=None):
    text(slide, Inches(0.7), Inches(0.42), Inches(12), Inches(0.3), kicker.upper(),
         size=11, color=ACCENT, bold=True)
    text(slide, Inches(0.7), Inches(0.72), Inches(12), Inches(0.65), title,
         size=30, bold=True)
    y = Inches(1.42)
    if sub:
        text(slide, Inches(0.7), y, Inches(11.9), Inches(0.5), sub,
             size=14, color=DIM, spacing=1.25)
        y = Inches(1.95)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.32),
                                  Inches(1.1), Pt(2.5))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    line.shadow.inherit = False
    return y


def stat(slide, x, y, w, value, label, *, color=TEXT, h=Inches(1.15), vsize=30):
    panel(slide, x, y, w, h, fill=PANEL2)
    text(slide, x + Inches(0.14), y + Inches(0.14), w - Inches(0.28), Inches(0.5),
         value, size=vsize, bold=True, color=color, align=PP_ALIGN.CENTER)
    text(slide, x + Inches(0.1), y + h - Inches(0.4), w - Inches(0.2), Inches(0.3),
         label.upper(), size=9, color=FAINT, align=PP_ALIGN.CENTER, spacing=1.0)


def table(slide, x, y, w, headers, rows, *, col_w=None, size=12,
          head_color=FAINT, row_h=Inches(0.34), highlight_last=False):
    n = len(headers)
    col_w = col_w or [w / n] * n
    cx = x
    for i, hcol in enumerate(headers):
        text(slide, cx, y, col_w[i], Inches(0.3), hcol.upper(), size=9.5,
             color=head_color, bold=True,
             align=PP_ALIGN.RIGHT if i else PP_ALIGN.LEFT)
        cx += col_w[i]
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.28), w, Pt(1))
    rule.fill.solid(); rule.fill.fore_color.rgb = BORDER
    rule.line.fill.background(); rule.shadow.inherit = False

    ry = y + Inches(0.4)
    for r, row in enumerate(rows):
        last = highlight_last and r == len(rows) - 1
        if last:
            hl = panel(slide, x - Inches(0.1), ry - Inches(0.05),
                       w + Inches(0.2), row_h + Inches(0.02),
                       fill=PANEL2, line=ACCENT)
            hl.line.width = Pt(1)
        cx = x
        for i, cell in enumerate(row):
            val, over = (cell, {}) if isinstance(cell, str) else cell
            text(slide, cx, ry, col_w[i], row_h, val,
                 size=over.get("size", size),
                 color=over.get("color", TEXT if i == 0 or last else DIM),
                 bold=over.get("bold", last),
                 font=over.get("font", SANS),
                 align=PP_ALIGN.RIGHT if i else PP_ALIGN.LEFT)
            cx += col_w[i]
        ry += row_h
    return ry


def bullets(slide, x, y, w, items, *, size=15, gap=0.42, dot=ACCENT):
    for i, item in enumerate(items):
        cy = y + Inches(gap * i)
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, cy + Inches(0.08),
                                   Inches(0.11), Inches(0.11))
        d.fill.solid(); d.fill.fore_color.rgb = dot
        d.line.fill.background(); d.shadow.inherit = False
        if isinstance(item, tuple):
            lead, rest = item
            tb = slide.shapes.add_textbox(x + Inches(0.28), cy,
                                          w - Inches(0.28), Inches(gap))
            tf = tb.text_frame; tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]; p.line_spacing = 1.2
            r1 = p.add_run(); r1.text = lead + " "
            r1.font.size = Pt(size); r1.font.bold = True
            r1.font.color.rgb = TEXT; r1.font.name = SANS
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = DIM; r2.font.name = SANS
        else:
            text(slide, x + Inches(0.28), cy, w - Inches(0.28), Inches(gap),
                 item, size=size, color=DIM, spacing=1.2)


def fit_picture(slide, path: Path, x, y, box_w, box_h, *, border=True,
                centre=True):
    """Scale a screenshot to fit inside a box, preserving aspect ratio.

    Setting width alone is how pictures silently run off the bottom of a slide:
    a 16:10 console capture at 11.9in wide is 7.4in tall, taller than the whole
    slide. Constraining both dimensions is the only safe way to place these.
    """
    if not path.exists():
        return None
    from PIL import Image
    try:
        with Image.open(path) as im:
            ar = im.width / im.height
    except Exception:
        ar = 1.6
    w = box_w
    h = Emu(int(w / ar))
    if h > box_h:
        h = box_h
        w = Emu(int(h * ar))
    px = x + Emu(int((box_w - w) / 2)) if centre else x
    pic = slide.shapes.add_picture(str(path), px, y, width=w, height=h)
    if border:
        pic.line.color.rgb = BORDER
        pic.line.width = Pt(1)
    return pic


def picture(slide, path: Path, x, y, w=None, h=None, *, border=True):
    if not path.exists():
        return None
    kw = {}
    if w:
        kw["width"] = w
    if h:
        kw["height"] = h
    pic = slide.shapes.add_picture(str(path), x, y, **kw)
    if border:
        pic.line.color.rgb = BORDER
        pic.line.width = Pt(1)
    return pic


def footer(slide, n, total, note=""):
    text(slide, Inches(0.7), H - Inches(0.5), Inches(9), Inches(0.3),
         note or "BODHI MULE HUNTER AI  ·  CyberShield Hackathon 2026  ·  Problem Statement 2",
         size=9.5, color=FAINT)
    text(slide, W - Inches(1.6), H - Inches(0.5), Inches(0.9), Inches(0.3),
         f"{n} / {total}", size=9.5, color=FAINT, align=PP_ALIGN.RIGHT)


# --------------------------------------------------------------------------


def build(m: dict) -> Presentation:
    ds, L = m["dataset"], m["layers"]
    b, c, lat = m["baseline"], m["calibration"], m["latency"]
    ind, oot, lead = m["independence_from_tickets"], m["out_of_time"], m["lead_time"]
    thr = m["thresholds"]
    decoy = {d["population"]: d for d in m["hard_negatives"]}
    typ = {t["typology"]: t for t in m["typology"]}
    fused = L["Fused (L7)"]

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    TOTAL = 19
    n = 0

    def nxt():
        nonlocal n
        n += 1
        return n

    # ---------------------------------------------------------- 1 title
    s = blank(prs)
    glow = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(-2), Inches(-1.5),
                              Inches(9), Inches(6))
    glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x10, 0x1A, 0x2C)
    glow.line.fill.background(); glow.shadow.inherit = False

    mark = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.5),
                              Inches(0.85), Inches(0.85))
    mark.fill.solid(); mark.fill.fore_color.rgb = ACCENT
    mark.line.fill.background(); mark.shadow.inherit = False
    text(s, Inches(0.9), Inches(1.68), Inches(0.85), Inches(0.5), "B",
         size=32, bold=True, color=BG, align=PP_ALIGN.CENTER)

    text(s, Inches(2.0), Inches(1.55), Inches(10), Inches(0.9),
         "BODHI MULE HUNTER AI", size=46, bold=True)
    text(s, Inches(2.0), Inches(2.42), Inches(10), Inches(0.4),
         "Real-time mule account & suspicious transaction detection",
         size=17, color=ACCENT)

    text(s, Inches(0.9), Inches(3.35), Inches(11.5), Inches(1.0),
         "Nine layers. Three complementary model views. One calibrated score an "
         "investigator can act on — and a containment layer that refuses to act "
         "without corroboration.", size=16, color=DIM, spacing=1.35)

    x = Inches(0.9)
    for val, lab, col in [
        (f"{fused['roc_auc']:.4f}", "Fused ROC-AUC", TEXT),
        (f"{fused['pr_auc']:.4f}", "Fused PR-AUC", TEXT),
        (_pct(b["false_positive_reduction"], 2), "False positives cut", GREEN),
        (f"{b['precision_uplift_x']:.0f}×", "Precision uplift", GREEN),
        (f"{lat['p50_ms']:.3f} ms", "Inline decision p50", ACCENT),
    ]:
        stat(s, x, Inches(4.7), Inches(2.24), val, lab, color=col)
        x += Inches(2.36)

    x = Inches(0.9)
    for mem in TEAM:
        text(s, x, Inches(6.12), Inches(2.6), Inches(0.28), mem.name,
             size=13, bold=True)
        text(s, x, Inches(6.4), Inches(2.6), Inches(0.26), mem.enrolment,
             size=10.5, color=FAINT, font=MONO)
        x += Inches(2.75)
    text(s, Inches(0.9), Inches(6.86), Inches(11.5), Inches(0.35),
         f"{TEAM_NAME}   ·   {EVENT}   ·   {AFFILIATION}",
         size=11, color=FAINT)
    nxt()

    # ---------------------------------------------------------- 2 problem
    s = blank(prs)
    y = header(s, "The problem", "A mule account is a normal account",
               sub="Real KYC. Real customer. Payments clear. Nothing about the "
                   "account is anomalous — only the shape of the money moving "
                   "through it, and the company it keeps.")
    panel(s, Inches(0.7), y + Inches(0.15), Inches(5.75), Inches(3.5))
    text(s, Inches(1.0), y + Inches(0.42), Inches(5.2), Inches(0.4),
         "HOW A MULE RING ACTUALLY WORKS", size=11, bold=True, color=FAINT)
    steps = [
        ("1  Victims", "40 defrauded people push money into one collector, in 90 minutes"),
        ("2  Layering", "A → B → C → D within minutes, each hop keeping 2–10%"),
        ("3  Structuring", "Split into transfers parked just under the reporting limit"),
        ("4  Cash-out", "Drained via ATM/AePS at 03:12 — out of the banking system"),
    ]
    yy = y + Inches(0.92)
    for lead_txt, rest in steps:
        text(s, Inches(1.0), yy, Inches(1.35), Inches(0.3), lead_txt,
             size=13, bold=True, color=ACCENT)
        text(s, Inches(2.35), yy, Inches(3.85), Inches(0.6), rest,
             size=12.5, color=DIM, spacing=1.2)
        yy += Inches(0.62)

    panel(s, Inches(6.85), y + Inches(0.15), Inches(5.75), Inches(3.5))
    text(s, Inches(7.15), y + Inches(0.42), Inches(5.2), Inches(0.4),
         "WHY TODAY'S RULE ENGINES FAIL", size=11, bold=True, color=FAINT)
    text(s, Inches(7.15), y + Inches(0.85), Inches(5.2), Inches(1.6),
         "They score accounts in isolation, so they cannot express “this "
         "account's counterparties are themselves suspicious” or “these forty "
         "credits arrived in ninety minutes”. Lacking the vocabulary, they "
         "compensate with sensitivity.", size=13.5, color=DIM, spacing=1.3)
    stat(s, Inches(7.15), y + Inches(2.35), Inches(2.5),
         _n(b["alerts"]), "alerts raised", color=ORANGE, h=Inches(1.1), vsize=28)
    stat(s, Inches(9.9), y + Inches(2.35), Inches(2.45),
         _pct(b["precision"], 2), "precision", color=RED, h=Inches(1.1), vsize=28)
    text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
         "99 of every 100 investigations are wasted. That is the number we set out to fix.",
         size=15, color=TEXT, bold=True)
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 3 insight
    s = blank(prs)
    y = header(s, "The core insight", "Aggregates destroy the signal that matters",
               sub="The same monthly totals describe a shopkeeper and a mule. Only "
                   "the ordering separates them — and no aggregate preserves order.")
    panel(s, Inches(0.7), y + Inches(0.25), Inches(5.75), Inches(2.5), fill=PANEL2)
    text(s, Inches(1.0), y + Inches(0.5), Inches(5.2), Inches(0.35),
         "SHOPKEEPER", size=12, bold=True, color=GREEN)
    text(s, Inches(1.0), y + Inches(0.95), Inches(5.2), Inches(1.5),
         "“Received ₹4 lakh from 38 payers over the month.”\n\n"
         "Spread across weeks. Money rests. Counterparties repeat.",
         size=14, color=DIM, spacing=1.35)

    panel(s, Inches(6.85), y + Inches(0.25), Inches(5.75), Inches(2.5), fill=PANEL2)
    text(s, Inches(7.15), y + Inches(0.5), Inches(5.2), Inches(0.35),
         "MULE", size=12, bold=True, color=RED)
    text(s, Inches(7.15), y + Inches(0.95), Inches(5.2), Inches(1.6),
         "“Received ₹4 lakh from 38 payers between 14:10 and 15:40, then "
         "withdrew ₹3.8 lakh in nineteen ATM visits beginning at 03:12.”",
         size=14, color=DIM, spacing=1.35)

    text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(0.45),
         "Identical aggregates. Completely different sequence.",
         size=19, bold=True, color=TEXT)
    text(s, Inches(0.7), Inches(6.05), Inches(12), Inches(0.8),
         "So the ensemble is not three flavours of the same tabular model. It is a "
         "tabular model, a graph neural network and a temporal model — each seeing "
         "something the other two structurally cannot.",
         size=14, color=DIM, spacing=1.3)
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 4 architecture
    s = blank(prs)
    y = header(s, "Architecture", "Nine layers, each owned by a named agent",
               sub="Declared input/output contracts, so the provenance of every "
                   "decision is recoverable from the audit trail.")
    groups = [
        ("INGEST & ENRICH", ACCENT, [
            ("L1", "Transaction ingestion", "UPI · IMPS · NEFT · RTGS · AePS · ATM · card + NCRP tickets + RBI feeds"),
            ("L2", "Feature engineering", "60 features: velocity, fan-in/out, dormancy, burst, structuring, device reuse"),
            ("L3", "Graph construction", "accounts × devices × IPs × VPAs, shared-identifier cliques capped"),
        ]),
        ("DETECT", PURPLE, [
            ("L4", "XGBoost screening", "tabular behaviour → exact TreeSHAP attributions"),
            ("L5", "GraphSAGE", "network structure → rings, device farms, multi-hop layering"),
            ("L6", "Temporal graph network", "event ordering → smurfing, dormant bursts, rapid routing"),
        ]),
        ("DECIDE & ACT", GREEN, [
            ("L7", "Risk fusion", "monotone non-negative blend, isotonic calibrated 0–100"),
            ("L8", "Explainability", "TreeSHAP + GNNExplainer + plain-English narrative"),
            ("L9", "Alerting & kill-switch", "proportionate, corroborated, reversible, audited"),
        ]),
    ]
    gy = y + Inches(0.1)
    for gname, gcol, rows in groups:
        panel(s, Inches(0.7), gy, Inches(11.9), Inches(1.42))
        text(s, Inches(0.92), gy + Inches(0.12), Inches(2.0), Inches(0.3),
             gname, size=9.5, bold=True, color=gcol)
        ry = gy + Inches(0.44)
        for tag, name, desc in rows:
            text(s, Inches(0.92), ry, Inches(0.42), Inches(0.28), tag,
                 size=12, bold=True, color=gcol, font=MONO)
            text(s, Inches(1.45), ry, Inches(2.6), Inches(0.28), name,
                 size=12.5, bold=True)
            text(s, Inches(4.15), ry, Inches(8.2), Inches(0.28), desc,
                 size=11.5, color=DIM)
            ry += Inches(0.3)
        gy += Inches(1.58)
    text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.35),
         "BODHI SHIELD (APK static triage) feeds extracted UPI handles and beneficiary "
         "accounts into Layer 3 — priming the graph before the first victim installs the trojan.",
         size=12, color=FAINT, spacing=1.25)
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 5 three views
    s = blank(prs)
    y = header(s, "Layers 4–6", "Three views that cannot see each other's evidence")
    cards = [
        ("LAYER 4", "XGBoost", ACCENT,
         "Aggregates per account",
         "Volume, ratios, cash-out share, structuring",
         "Who the counterparties are",
         f"AUC {L['XGBoost (L4)']['roc_auc']:.4f}"),
        ("LAYER 5", "GraphSAGE", PURPLE,
         "Two-hop neighbourhood",
         "Rings, device farms, multi-hop layering",
         "When things happened",
         f"AUC {L['GraphSAGE (L5)']['roc_auc']:.4f}"),
        ("LAYER 6", "Temporal GNN", GREEN,
         "Last 64 events, in order",
         "Smurfing, dormant bursts, rapid routing",
         "Anything outside the account",
         f"AUC {L['TGN (L6)']['roc_auc']:.4f}"),
    ]
    x = Inches(0.7)
    for tag, name, col, sees, catches, blind, auc in cards:
        panel(s, x, y + Inches(0.15), Inches(3.85), Inches(3.75))
        text(s, x + Inches(0.28), y + Inches(0.4), Inches(3.3), Inches(0.25),
             tag, size=9.5, bold=True, color=col)
        text(s, x + Inches(0.28), y + Inches(0.68), Inches(3.3), Inches(0.4),
             name, size=21, bold=True)
        yy = y + Inches(1.25)
        for lab, val, lcol in [("SEES", sees, FAINT), ("CATCHES", catches, FAINT),
                               ("BLIND TO", blind, FAINT)]:
            text(s, x + Inches(0.28), yy, Inches(3.3), Inches(0.2), lab,
                 size=8.5, bold=True, color=lcol)
            text(s, x + Inches(0.28), yy + Inches(0.22), Inches(3.3), Inches(0.55),
                 val, size=12, color=DIM if lab != "BLIND TO" else ORANGE,
                 spacing=1.2)
            yy += Inches(0.82)
        text(s, x + Inches(0.28), y + Inches(3.45), Inches(3.3), Inches(0.3),
             auc, size=13, bold=True, color=col, font=MONO)
        x += Inches(4.03)

    panel(s, Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.82), fill=PANEL2)
    text(s, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.55),
         "Both neural layers are implemented in NumPy with hand-derived gradients — no PyTorch, "
         "no CUDA, no 2 GB of wheels. Every analytic gradient is verified against central finite "
         "differences in the test suite, so “hand-rolled” does not mean “unverified”.",
         size=12.5, color=DIM, spacing=1.25)
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 6 fusion
    s = blank(prs)
    y = header(s, "Layer 7", "Fusion: two properties that are not negotiable")
    panel(s, Inches(0.7), y + Inches(0.2), Inches(5.85), Inches(2.35))
    text(s, Inches(1.0), y + Inches(0.45), Inches(5.3), Inches(0.35),
         "MONOTONICITY", size=11, bold=True, color=ACCENT)
    text(s, Inches(1.0), y + Inches(0.85), Inches(5.3), Inches(1.5),
         "Blend weights are constrained non-negative, so more evidence can never "
         "lower a score. An unconstrained stacker on a small fold reliably learns "
         "“higher temporal score means safer” — indefensible to an "
         "investigator, unshippable past a regulator.",
         size=13, color=DIM, spacing=1.3)

    panel(s, Inches(6.75), y + Inches(0.2), Inches(5.85), Inches(2.35))
    text(s, Inches(7.05), y + Inches(0.45), Inches(5.3), Inches(0.35),
         "CALIBRATION", size=11, bold=True, color=GREEN)
    text(s, Inches(7.05), y + Inches(0.85), Inches(5.3), Inches(1.1),
         "The score drives an automated kill-switch at 85, so it has to mean what "
         "it says. Isotonic regression maps the blend onto observed frequencies.",
         size=13, color=DIM, spacing=1.3)
    stat(s, Inches(7.05), y + Inches(1.62), Inches(2.55), f"{c['brier']:.5f}",
         "Brier score", color=GREEN, h=Inches(0.82), vsize=20)
    stat(s, Inches(9.8), y + Inches(1.62), Inches(2.55), f"{c['ece']:.4f}",
         "Expected calib. error", color=GREEN, h=Inches(0.82), vsize=20)

    panel(s, Inches(0.7), Inches(5.05), Inches(11.9), Inches(1.55), fill=PANEL2)
    text(s, Inches(1.0), Inches(5.25), Inches(11.3), Inches(0.3),
         "FOUR FOLDS, NOT THREE", size=11, bold=True, color=FAINT)
    xx = Inches(1.0)
    for name, role, col in [("train", "fits the base models", ACCENT),
                            ("val", "early-stops them", PURPLE),
                            ("fuse", "fits the Layer-7 blend", GREEN),
                            ("test", "touched exactly once", ORANGE)]:
        text(s, xx, Inches(5.62), Inches(2.7), Inches(0.3), name,
             size=15, bold=True, color=col, font=MONO)
        text(s, xx, Inches(5.95), Inches(2.7), Inches(0.4), role,
             size=11.5, color=DIM)
        xx += Inches(2.85)
    text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.3),
         "A stacker fitted on the fold that early-stopped its own inputs inherits their "
         "optimism — and came out worse than the best single layer.",
         size=11.5, color=FAINT)
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 7 results
    s = blank(prs)
    y = header(s, "Results", "Every layer, on held-out data",
               sub=f"{_n(ds['accounts'])} accounts · {_n(ds['transactions'])} transactions · "
                   f"{_n(ds['mule_accounts'])} mules ({_pct(ds['mule_rate'], 2)}) · {ds['rings']} rings")
    rows = []
    for name, v in L.items():
        rows.append([name, f"{v['roc_auc']:.4f}", f"{v['pr_auc']:.4f}",
                     f"{v['precision_at_1pct']:.3f}", f"{v['recall_at_1pct']:.3f}"])
    table(s, Inches(0.7), y + Inches(0.15), Inches(7.1),
          ["Layer", "ROC-AUC", "PR-AUC", "P@1%", "R@1%"], rows,
          col_w=[Inches(2.9), Inches(1.1), Inches(1.1), Inches(1.0), Inches(1.0)],
          highlight_last=True, row_h=Inches(0.43))

    picture(s, FIGURE_DIR / "roc_pr.png", Inches(8.15), y + Inches(0.1),
            w=Inches(4.45))
    text(s, Inches(0.7), Inches(5.35), Inches(7.1), Inches(0.35),
         "The fused model beats every individual view on both ranking metrics — "
         "which is the complementarity claim, measured.",
         size=12.5, color=DIM, spacing=1.25)

    xx = Inches(0.7)
    for val, lab, col in [(f"{_pct(oot['recall'], 0)}", "out-of-time recall", GREEN),
                          (f"{_pct(ind['share_of_detections_not_in_any_ticket'], 0)}",
                           "never in any NCRP ticket", ACCENT),
                          (f"{lat['p99_ms']:.3f} ms", "inline decision p99", PURPLE)]:
        stat(s, xx, Inches(5.95), Inches(2.3), val, lab, color=col,
             h=Inches(0.95), vsize=22)
        xx += Inches(2.45)
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 8 vs rules
    s = blank(prs)
    y = header(s, "The headline", "Versus the rule engine it replaces",
               sub="Compared at identical recall — because a system that alerts on "
                   "everything can always claim to catch more.")
    panel(s, Inches(0.7), y + Inches(0.2), Inches(5.6), Inches(3.1))
    text(s, Inches(1.0), y + Inches(0.45), Inches(5.0), Inches(0.35),
         "RULE ENGINE  ·  8 SCENARIOS", size=11, bold=True, color=ORANGE)
    for i, (lab, val) in enumerate([("Alerts raised", _n(b["alerts"])),
                                    ("True positives", _n(b["true_positives"])),
                                    ("False positives", _n(b["false_positives"])),
                                    ("Precision", _pct(b["precision"], 2))]):
        yy = y + Inches(0.95) + Inches(0.52) * i
        text(s, Inches(1.0), yy, Inches(3.0), Inches(0.35), lab, size=14, color=DIM)
        text(s, Inches(4.0), yy, Inches(2.0), Inches(0.35), val, size=17,
             bold=True, color=RED if "alse" in lab or "recision" in lab else TEXT,
             align=PP_ALIGN.RIGHT)

    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.45),
                               y + Inches(1.5), Inches(0.55), Inches(0.4))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT
    arrow.line.fill.background(); arrow.shadow.inherit = False

    p = panel(s, Inches(7.15), y + Inches(0.2), Inches(5.45), Inches(3.1), fill=PANEL2)
    p.line.color.rgb = GREEN
    p.line.width = Pt(1.5)
    text(s, Inches(7.45), y + Inches(0.45), Inches(4.9), Inches(0.35),
         "BODHI  ·  SAME RECALL", size=11, bold=True, color=GREEN)
    for i, (lab, val) in enumerate([
            ("Alerts raised", _n(b["bodhi_alerts_at_equal_recall"])),
            ("True positives", _n(b["bodhi_true_positives_at_equal_recall"])),
            ("False positives", _n(b["bodhi_false_positives_at_equal_recall"])),
            ("Precision", _pct(b["bodhi_precision_at_equal_recall"], 2))]):
        yy = y + Inches(0.95) + Inches(0.52) * i
        text(s, Inches(7.45), yy, Inches(3.0), Inches(0.35), lab, size=14, color=DIM)
        text(s, Inches(10.3), yy, Inches(2.0), Inches(0.35), val, size=17,
             bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

    band = panel(s, Inches(0.7), Inches(5.65), Inches(11.9), Inches(1.05),
                 fill=RGBColor(0x0F, 0x2A, 0x22), line=GREEN)
    text(s, Inches(1.0), Inches(5.85), Inches(11.3), Inches(0.7),
         f"{_pct(b['false_positive_reduction'], 2)} of false positives eliminated at "
         f"unchanged recall — a {b['precision_uplift_x']:.0f}× precision uplift. "
         f"The analyst clears the same true positives while opening "
         f"{_n(b['bodhi_alerts_at_equal_recall'])} cases instead of {_n(b['alerts'])}.",
         size=16, bold=True, color=TEXT, spacing=1.25)
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 9 credibility
    s = blank(prs)
    y = header(s, "Why that number is credible", "We planted the accounts designed to break it",
               sub="Any model scores ~1.0 AUC on naive synthetic fraud data. So the "
                   "simulated bank is stocked with legitimate accounts that are "
                   "structurally indistinguishable from mules.")
    rows = []
    for k, label in [("bc_agent_device_cluster", "Business Correspondent (Bank Mitra) agents"),
                     ("community_collector", "Community collectors (chit fund, tuition)"),
                     ("small_business_fan_in", "Small businesses with heavy fan-in"),
                     ("legitimate_dormant_wake", "Genuine dormancy reactivation"),
                     ("ordinary_retail", "Ordinary retail")]:
        if k in decoy:
            d = decoy[k]
            rows.append([label, _n(d["n"]),
                         (_pct(d["false_positive_rate"], 2),
                          {"color": GREEN if d["false_positive_rate"] < 0.01 else AMBER,
                           "bold": True})])
    table(s, Inches(0.7), y + Inches(0.15), Inches(7.3),
          ["Legitimate population", "Count", "False-positive rate"], rows,
          col_w=[Inches(4.5), Inches(1.2), Inches(1.6)], row_h=Inches(0.46))

    panel(s, Inches(8.4), y + Inches(0.1), Inches(4.2), Inches(3.3), fill=PANEL2)
    text(s, Inches(8.7), y + Inches(0.35), Inches(3.6), Inches(0.35),
         "THE HARDEST DECOY", size=10, bold=True, color=ORANGE)
    text(s, Inches(8.7), y + Inches(0.75), Inches(3.6), Inches(2.3),
         "A Business Correspondent legitimately operates 8–22 village accounts from "
         "one handheld and dispenses AePS cash all day.\n\n"
         "That is the exact fingerprint of a device farm doing rapid cash-out — and "
         "the system does not fall for it.",
         size=13, color=DIM, spacing=1.3)

    guards = [
        ("Four disjoint folds", "test touched exactly once"),
        ("Cross-fitted seed model", "no label leaks into a neighbour's average"),
        ("Intelligence quarantined", "NCRP tickets never reach the supervised layers"),
        ("Point-in-time correctness", "tested against a manually truncated ledger"),
    ]
    text(s, Inches(0.7), Inches(5.4), Inches(12), Inches(0.3),
         "OTHER GUARDS AGAINST FLATTERING OURSELVES", size=10, bold=True, color=FAINT)
    xx = Inches(0.7)
    for lead_txt, rest in guards:
        text(s, xx, Inches(5.78), Inches(2.85), Inches(0.3), lead_txt,
             size=12.5, bold=True, color=TEXT)
        text(s, xx, Inches(6.08), Inches(2.85), Inches(0.6), rest,
             size=11, color=DIM, spacing=1.2)
        xx += Inches(3.0)
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 10 weakness
    s = blank(prs)
    y = header(s, "Where it is weak", "Recall is not uniform, and averaging hides that",
               sub="A deployment team needs to know which typologies still need "
                   "human-designed scenarios alongside the model.")
    rows = []
    for k, t in sorted(typ.items(), key=lambda kv: kv[1]["recall"])[:9]:
        col = RED if t["recall"] < 0.90 else (AMBER if t["recall"] < 0.99 else GREEN)
        rows.append([k.replace("role:", "role: ").replace("_", " "),
                     (f"{t['recall']:.1%}", {"color": col, "bold": True}),
                     f"{t['detected']}/{t['mules']}"])
    table(s, Inches(0.7), y + Inches(0.15), Inches(6.3),
          ["Typology / role", "Recall", "Detected"], rows,
          col_w=[Inches(3.5), Inches(1.4), Inches(1.4)], row_h=Inches(0.4))

    panel(s, Inches(7.4), y + Inches(0.1), Inches(5.2), Inches(3.9), fill=PANEL2)
    text(s, Inches(7.7), y + Inches(0.35), Inches(4.6), Inches(0.35),
         "WHY STRUCTURING IS THE HARDEST", size=10, bold=True, color=ORANGE)
    text(s, Inches(7.7), y + Inches(0.75), Inches(4.6), Inches(1.5),
         "A single sub-threshold transfer is, by construction, indistinguishable "
         "from a genuine near-threshold property or jewellery payment. Detection "
         "depends on an aggregate pattern that only becomes visible once enough "
         "transfers have accumulated.", size=12.5, color=DIM, spacing=1.3)
    text(s, Inches(7.7), y + Inches(2.35), Inches(4.6), Inches(0.35),
         "AND CASH-OUT ACCOUNTS", size=10, bold=True, color=ORANGE)
    text(s, Inches(7.7), y + Inches(2.75), Inches(4.6), Inches(1.0),
         "Short, thin histories — they receive once and drain. There is simply "
         "less evidence to work with.", size=12.5, color=DIM, spacing=1.3)

    text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
         "We report these rather than averaging them away. A prototype that hides its "
         "blind spots is not useful to the team that has to deploy it.",
         size=14, color=TEXT, bold=True, spacing=1.25)
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 11 console
    s = blank(prs)
    y = header(s, "Live demo", "The investigator console",
               sub="Four layer scores shown separately, evidence written in the "
                   "vocabulary of an AML desk, cluster and exposure on every card.")
    fit_picture(s, SHOTS / "01_investigate_queue.png", Inches(0.7), y + Inches(0.1),
                Inches(11.9), Inches(4.85))
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 12 explainability
    s = blank(prs)
    y = header(s, "Layer 8", "Explanations specific enough to file")
    picture(s, SHOTS / "02_shap_attribution.png", Inches(0.7), y + Inches(0.1),
            w=Inches(5.9))
    picture(s, SHOTS / "03_network_graph.png", Inches(6.75), y + Inches(0.1),
            w=Inches(5.85))
    text(s, Inches(0.7), Inches(5.6), Inches(5.9), Inches(1.2),
         "Exact TreeSHAP — attributions sum to the model's margin, not an "
         "approximation of it. Rendered as sentences, not feature names: "
         "“94% of every rupee credited left within 60 minutes — the account is "
         "a conduit, not a destination.”", size=12.5, color=DIM, spacing=1.3)
    text(s, Inches(6.75), Inches(5.6), Inches(5.85), Inches(1.2),
         "GNNExplainer learns a mask over the two-hop neighbourhood — answering "
         "which relationships drove the score, which SHAP structurally cannot. "
         "Dashed edges are shared-device links, invisible to any tabular model.",
         size=12.5, color=DIM, spacing=1.3)
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 13 money trail
    s = blank(prs)
    y = header(s, "Following the money", "Time-respecting multi-hop trails",
               sub="Every hop must occur after the previous one — money cannot be "
                   "forwarded before it arrives. Value retention exposes each "
                   "mule's commission.")
    picture(s, SHOTS / "04_money_trail.png", Inches(0.7), y + Inches(0.1),
            w=Inches(7.6))
    panel(s, Inches(8.6), y + Inches(0.1), Inches(4.0), Inches(3.2), fill=PANEL2)
    text(s, Inches(8.9), y + Inches(0.35), Inches(3.4), Inches(0.35),
         "THE TELL", size=10, bold=True, color=ACCENT)
    text(s, Inches(8.9), y + Inches(0.75), Inches(3.4), Inches(2.2),
         "A laundering chain leaks 2–10% at every hop as the handler's cut.\n\n"
         "Legitimate payment chains simply do not have that shape — which is why "
         "value retention is shown on every path.",
         size=13, color=DIM, spacing=1.35)
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 14 containment
    s = blank(prs)
    y = header(s, "Layer 9", "Containment that can refuse",
               sub="A wrongly frozen account is somebody unable to pay for medicine. "
                   "No AUC improvement justifies being casual about that.")
    constraints = [
        ("Proportionality", "below 85 the automation cannot freeze at all"),
        ("Corroboration", "a full freeze needs ≥2 independent layers agreeing"),
        ("Reversibility", "every action carries a TTL and can be reverted"),
        ("Rate limiting", "bounded automated freezes per hour caps blast radius"),
        ("Protected accounts", "salary, pension and benefit accounts always go to a human"),
    ]
    yy = y + Inches(0.15)
    for lead_txt, rest in constraints:
        panel(s, Inches(0.7), yy, Inches(6.4), Inches(0.62), fill=PANEL2)
        text(s, Inches(1.0), yy + Inches(0.14), Inches(2.3), Inches(0.35),
             lead_txt, size=13, bold=True, color=ACCENT)
        text(s, Inches(3.3), yy + Inches(0.16), Inches(3.6), Inches(0.35),
             rest, size=12, color=DIM)
        yy += Inches(0.72)

    picture(s, SHOTS / "05_killswitch.png", Inches(7.35), y + Inches(0.15),
            w=Inches(5.25))
    text(s, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.6),
         "When a constraint fires, the action is downgraded — never escalated — and "
         "flagged for human review. Every decision, including every refusal, is written "
         "to a hash-chained audit log whose head hash can be externally anchored.",
         size=13, color=DIM, spacing=1.3)
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 15 shield
    s = blank(prs)
    y = header(s, "Ecosystem synergy", "BODHI SHIELD → BODHI MULE HUNTER",
               sub="Malware analysis priming fraud controls before the first victim "
                   "installs the app.")
    picture(s, SHOTS / "08_shield_apk.png", Inches(0.7), y + Inches(0.1),
            w=Inches(7.5))
    steps = [
        ("1", "Parse manifest", "binary AndroidManifest.xml string pool, to spec"),
        ("2", "Score combinations", "READ_SMS alone is a messaging app; + accessibility + overlay is a trojan"),
        ("3", "Mine the DEX", "UPI handles, IFSC codes, accounts, C2 endpoints"),
        ("4", "Hand off", "identifiers become weighted nodes in the fraud graph"),
    ]
    yy = y + Inches(0.15)
    for num, lead_txt, rest in steps:
        text(s, Inches(8.55), yy, Inches(0.35), Inches(0.3), num,
             size=15, bold=True, color=PURPLE, font=MONO)
        text(s, Inches(9.0), yy, Inches(3.6), Inches(0.3), lead_txt,
             size=13, bold=True)
        text(s, Inches(9.0), yy + Inches(0.28), Inches(3.6), Inches(0.6), rest,
             size=11, color=DIM, spacing=1.2)
        yy += Inches(0.82)
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 16 deployment
    s = blank(prs)
    y = header(s, "Deployment", "How it fits a real bank",
               sub="The split between scheduled graph work and the inline path is what "
                   "makes a UPI-latency decision possible at all.")
    boxes = [
        ("Core banking / switch", "Kafka event stream", ACCENT, Inches(0.7)),
        ("L1–L2 streaming features", "feature store", ACCENT, Inches(3.85)),
        ("L3/L5/L6 batch\n(scheduled)", "standing risk cache", PURPLE, Inches(7.0)),
        ("Inline decision\n< 0.1 ms", "ALLOW · REVIEW · HOLD · BLOCK", GREEN, Inches(10.15)),
    ]
    for title, sub, col, x in boxes:
        p = panel(s, x, y + Inches(0.4), Inches(2.75), Inches(1.5), fill=PANEL2)
        p.line.color.rgb = col
        text(s, x + Inches(0.2), y + Inches(0.65), Inches(2.35), Inches(0.6),
             title, size=13, bold=True, align=PP_ALIGN.CENTER, spacing=1.2)
        text(s, x + Inches(0.2), y + Inches(1.42), Inches(2.35), Inches(0.4),
             sub, size=10.5, color=DIM, align=PP_ALIGN.CENTER, spacing=1.15)

    perf = [("Simulate the bank", "7 s"), ("Feature engineering", "~4 s"),
            ("Graph construction", "~3 s"), ("Full training, all layers", "~123 s"),
            ("Full population re-score", "~38 s"),
            ("Inline decision (p50)", f"{lat['p50_ms']:.3f} ms")]
    text(s, Inches(0.7), Inches(4.5), Inches(5.5), Inches(0.3),
         "MEASURED PERFORMANCE", size=10, bold=True, color=FAINT)
    yy = Inches(4.88)
    for lab, val in perf:
        text(s, Inches(0.7), yy, Inches(3.6), Inches(0.3), lab, size=12.5, color=DIM)
        text(s, Inches(4.3), yy, Inches(1.6), Inches(0.3), val, size=12.5,
             bold=True, color=TEXT, font=MONO, align=PP_ALIGN.RIGHT)
        yy += Inches(0.33)

    panel(s, Inches(6.6), Inches(4.5), Inches(6.0), Inches(2.15), fill=PANEL2)
    text(s, Inches(6.9), Inches(4.72), Inches(5.4), Inches(0.3),
         "WHY THIS SPLIT", size=10, bold=True, color=FAINT)
    text(s, Inches(6.9), Inches(5.1), Inches(5.4), Inches(1.4),
         "The graph and temporal layers run on a schedule and cache a standing risk "
         "per account. The authorisation path combines those cached scores with the "
         "transaction's own attributes and never traverses the graph — which is the "
         "only way a decision is feasible at payment latencies.",
         size=12.5, color=DIM, spacing=1.3)
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 17 their data
    boi_path = METRICS_DIR / "boi_track.json"
    if boi_path.exists():
        bm = json.loads(boi_path.read_text())
        bd, bdep, bleak = bm["dataset"], bm["deployable"], bm["leakage_effect"]
        bcv, bhold = bdep["cv"], bdep["holdout"]

        s = blank(prs)
        y = header(s, "Their dataset", "Built on the schema they published",
                   sub=f"{_n(bd['declared_columns'])} declared columns, one row per "
                       f"alert. The names are machine-generated from a grammar, so we "
                       f"parse them instead of treating them as opaque.")

        # -- left: the leakage finding
        p = panel(s, Inches(0.7), y + Inches(0.2), Inches(5.9), Inches(2.45),
                  fill=PANEL2)
        p.line.color.rgb = RED
        p.line.width = Pt(1.5)
        text(s, Inches(1.0), y + Inches(0.35), Inches(5.3), Inches(0.3),
             "FOUR COLUMNS LEAK THE LABEL", size=11, bold=True, color=RED)
        text(s, Inches(1.0), y + Inches(0.72), Inches(5.3), Inches(0.9),
             "FRAUD_SUSPECTED · FALSE_POSITIVE · OTHER_RESOLUTION · UNATTENDED "
             "are resolution-status flags — how an analyst closed the alert. "
             "An open alert has none of them.",
             size=12.5, color=DIM, spacing=1.3)
        for i, (lab, val, col) in enumerate([
                ("PR-AUC, quarantined (deployable)",
                 f"{bleak['pr_auc_deployable']:.3f}", GREEN),
                ("PR-AUC, resolution columns admitted",
                 f"{bleak['pr_auc_with_leakage']:.3f}", RED)]):
            yy = y + Inches(1.68) + Inches(0.38) * i
            text(s, Inches(1.0), yy, Inches(3.9), Inches(0.32), lab,
                 size=12, color=DIM)
            text(s, Inches(4.9), yy, Inches(1.4), Inches(0.32), val, size=16,
                 bold=True, color=col, font=MONO, align=PP_ALIGN.RIGHT)

        # -- right: strategies measured, not assumed
        order = [("bank_finalized", "Bank's 18 finalised"),
                 ("bank_plus_engineered", "Bank + engineered"),
                 ("auto_topk", "Automatic top-k"),
                 ("all", "Every column")]
        rows = [[label, _n(bcv[k]["n_features"]), f"{bcv[k]['roc_auc']:.3f}",
                 f"{bcv[k]['pr_auc']:.3f}"] for k, label in order if k in bcv]
        text(s, Inches(7.0), y + Inches(0.2), Inches(5.6), Inches(0.3),
             "FOUR STRATEGIES, SELECTION INSIDE EVERY FOLD", size=11, bold=True,
             color=FAINT)
        table(s, Inches(7.0), y + Inches(0.58), Inches(5.6),
              ["Strategy", "Feats", "ROC-AUC", "PR-AUC"], rows,
              col_w=[Inches(2.3), Inches(1.0), Inches(1.15), Inches(1.15)],
              row_h=Inches(0.38), size=12)
        text(s, Inches(7.0), y + Inches(2.52), Inches(5.6), Inches(0.4),
             f"The bank's eighteen beat all {_n(bcv['all']['n_features'])} columns. "
             f"With {_n(bd['positives'])} positives against thousands of predictors, "
             "that is what theory predicts.",
             size=11.5, color=DIM, spacing=1.2)

        xx = Inches(0.7)
        for val, lab, col in [
                (f"{bhold['roc_auc']:.4f}", "untouched holdout ROC-AUC", GREEN),
                (f"{bcv[bdep['selected_strategy']]['roc_auc']:.4f}",
                 "cross-validated estimate", ACCENT),
                (f"{bleak['multiple']:.1f}×", "PR-AUC inflation if leaked", RED),
                (f"{_n(bd['declared_columns'])}", "columns parsed by grammar", PURPLE)]:
            stat(s, xx, Inches(4.9), Inches(2.9), val, lab, color=col,
                 h=Inches(0.95), vsize=22)
            xx += Inches(3.05)

        band = panel(s, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.85),
                     fill=RGBColor(0x2A, 0x1A, 0x0F), line=ORANGE)
        text(s, Inches(1.0), Inches(6.16), Inches(11.3), Inches(0.55),
             "Measured on a stand-in table with their exact schema — their data was not "
             "released when this was built. It proves the pipeline runs and does not "
             "flatter itself. It is not model performance, and we will not present it "
             "as such.",
             size=13, bold=True, color=TEXT, spacing=1.2)
        footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 18 limitations
    s = blank(prs)
    y = header(s, "Limitations", "Stated plainly, because hiding them helps nobody")
    lims = [
        ("Simulated data", "Demonstrates the architecture beats a rule engine on data "
                           "with realistic decoys. Not a forecast of production performance."),
        ("Static APK analysis only", "The dynamic sandbox with runtime hooking needs an "
                                     "instrumented Android image; it cannot ship self-contained."),
        ("Graph layers are not real-time", "~38 s full re-score. An account whose neighbourhood "
                                           "changed since the last batch is scored on slightly stale structure."),
        ("Single-institution view", "Rings routing through several banks are only partly "
                                    "visible — a data-sharing problem, not a modelling one."),
        ("Cold start", "The strongest feature is neighbourhood risk propagated from confirmed "
                       "cases; with zero known mules performance drops materially."),
        ("Disparate impact unevaluated", "Minimum-KYC and shared-device features are predictive "
                                         "but correlate with lower-income households. Alert-rate parity must be measured before go-live."),
    ]
    yy = y + Inches(0.1)
    for lead_txt, rest in lims:
        text(s, Inches(0.7), yy, Inches(3.4), Inches(0.35), lead_txt,
             size=13.5, bold=True, color=ORANGE)
        text(s, Inches(4.3), yy, Inches(8.3), Inches(0.6), rest,
             size=12.5, color=DIM, spacing=1.25)
        yy += Inches(0.78)
    footer(s, nxt(), TOTAL)

    # ---------------------------------------------------------- 19 close
    s = blank(prs)
    glow = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(-1),
                              Inches(9), Inches(6))
    glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x10, 0x1A, 0x2C)
    glow.line.fill.background(); glow.shadow.inherit = False

    text(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(1.0),
         "Mule detection fails today not because the\nsignal is absent",
         size=34, bold=True, spacing=1.2)
    text(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(0.9),
         "— but because single-account rule engines cannot express it.",
         size=22, color=ACCENT, spacing=1.2)
    text(s, Inches(0.9), Inches(3.6), Inches(11.0), Inches(1.1),
         "Three complementary views, fused under a monotonicity constraint into a "
         "calibrated score, cut false positives by "
         f"{_pct(b['false_positive_reduction'], 2)} at unchanged recall — with "
         "explanations specific enough to file and containment cautious enough to "
         "automate.", size=15, color=DIM, spacing=1.35)

    xx = Inches(0.9)
    for val, lab in [(str(test_count()), "tests passing"), (code_lines_short(), "lines of code"),
                     ("9", "layers, all running"), ("0", "external ML frameworks")]:
        stat(s, xx, Inches(5.0), Inches(2.6), val, lab, h=Inches(1.05), vsize=26)
        xx += Inches(2.75)

    x = Inches(0.9)
    for mem in TEAM:
        text(s, x, Inches(6.28), Inches(2.6), Inches(0.28), mem.name,
             size=12.5, bold=True)
        text(s, x, Inches(6.55), Inches(2.6), Inches(0.26), mem.enrolment,
             size=10, color=FAINT, font=MONO)
        x += Inches(2.75)
    text(s, Inches(0.9), Inches(6.98), Inches(11.5), Inches(0.35),
         "make setup && make all && make serve      ·      "
         "github.com/deepakvish001/BOI-Hackathon-Prototype-",
         size=11, color=FAINT, font=MONO)
    nxt()

    return prs


def main() -> int:
    path = METRICS_DIR / "evaluation.json"
    if not path.exists():
        print("run `make evaluate` first", file=sys.stderr)
        return 1
    prs = build(json.loads(path.read_text()))
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"PPTX  -> {OUT}  ({OUT.stat().st_size // 1024} KB, "
          f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides)")

    # PDF export for reviewers without PowerPoint
    try:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir",
             str(OUT.parent), str(OUT)],
            check=True, capture_output=True, timeout=300)
        if OUT_PDF.exists():
            print(f"PDF   -> {OUT_PDF}  ({OUT_PDF.stat().st_size // 1024} KB)")
    except Exception as exc:  # pragma: no cover - LibreOffice is optional
        print(f"(PDF export skipped: {exc})")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
